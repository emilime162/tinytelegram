#!/usr/bin/env python3
"""Generate TinyTelegram CS6650 Presentation — Northeastern Black/Red Theme"""

import os, sys
from io import BytesIO
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════════════════════════════════════════
#  CONSTANTS
# ═══════════════════════════════════════════════════════════
SW, SH = 13.333, 7.5
TOTAL_SLIDES = 10

# Northeastern brand colors (PPTX)
RED     = RGBColor(0xCC, 0x00, 0x00)
DKRED   = RGBColor(0x8B, 0x00, 0x00)
BLACK   = RGBColor(0x00, 0x00, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DKGRAY  = RGBColor(0x33, 0x33, 0x33)
MDGRAY  = RGBColor(0x88, 0x88, 0x88)
LTGRAY  = RGBColor(0xE0, 0xE0, 0xE0)
VLIGHT  = RGBColor(0xF5, 0xF5, 0xF5)
GREEN   = RGBColor(0x2E, 0x7D, 0x32)

# Matplotlib palette
M = dict(red='#CC0000', dkred='#8B0000', blk='#1a1a1a',
         dk='#333333', md='#888888', lt='#E8E8E8', wh='#FFFFFF',
         grn='#2E7D32')

plt.rcParams.update({
    'font.family': 'sans-serif',
    'font.sans-serif': ['Segoe UI', 'Calibri', 'Arial'],
    'axes.spines.top': False, 'axes.spines.right': False,
    'axes.labelcolor': '#333', 'xtick.color': '#555', 'ytick.color': '#555',
})

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   '..', 'TinyTelegram_CS6650.pptx')

# ═══════════════════════════════════════════════════════════
#  PPTX HELPERS
# ═══════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BL = prs.slide_layouts[6]

def setbg(sl, c):
    f = sl.background.fill; f.solid(); f.fore_color.rgb = c

def box(sl, l, t, w, h, fill=None, line=None, rd=False):
    s = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rd else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = Pt(1.2)
    else: s.line.fill.background()
    return s

def tx(sl, l, t, w, h, text, sz=18, c=BLACK, b=False,
       al=PP_ALIGN.LEFT, fn='Segoe UI', anc=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anc
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
    p.font.bold = b; p.font.name = fn; p.alignment = al
    return tb

def ml(sl, l, t, w, h, lines):
    """lines: [(text, size, color, bold, align), ...]"""
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, sz, c, b, a) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.bold = b; p.font.name = 'Segoe UI'; p.alignment = a
        p.space_before = Pt(0); p.space_after = Pt(1)
    return tb

def hdr(sl, title, sub=None):
    box(sl, 0, 0, SW, 0.1, fill=RED)
    tx(sl, 0.8, 0.28, 10, 0.5, title, sz=30, b=True)
    if sub: tx(sl, 0.8, 0.78, 10, 0.35, sub, sz=13, c=MDGRAY)
    box(sl, 0.8, 1.18, SW - 1.6, 0.012, fill=LTGRAY)

def ftr(sl, n):
    tx(sl, 0.8, SH - 0.38, 3, 0.25, 'TinyTelegram  |  CS6650', sz=9, c=MDGRAY)
    tx(sl, SW - 1.5, SH - 0.38, 1, 0.25, f'{n}/{TOTAL_SLIDES}',
       sz=9, c=MDGRAY, al=PP_ALIGN.RIGHT)

def addimg(sl, buf, l, t, w=None, h=None):
    kw = {}
    if w: kw['width'] = Inches(w)
    if h: kw['height'] = Inches(h)
    sl.shapes.add_picture(buf, Inches(l), Inches(t), **kw)

def tobuf(fig, dpi=200):
    b = BytesIO()
    fig.savefig(b, format='png', dpi=dpi, bbox_inches='tight',
                facecolor='white', edgecolor='none', pad_inches=0.15)
    plt.close(fig); b.seek(0); return b


# ═══════════════════════════════════════════════════════════
#  CHART: ARCHITECTURE DIAGRAM
# ═══════════════════════════════════════════════════════════
def make_arch():
    fig, ax = plt.subplots(figsize=(11, 5))
    ax.set_xlim(-0.2, 11.2); ax.set_ylim(-0.2, 5.5)
    ax.axis('off'); fig.patch.set_facecolor('white')

    def bx(cx, cy, w, h, color, label, sub=''):
        p = FancyBboxPatch((cx - w/2, cy - h/2), w, h,
                           boxstyle="round,pad=0.08",
                           facecolor=color, edgecolor='#444', lw=1.3,
                           zorder=3)
        ax.add_patch(p)
        ax.text(cx, cy + 0.02, label, ha='center', va='center',
                fontsize=10.5, fontweight='bold', color='white', zorder=4)
        if sub:
            ax.text(cx, cy - h/2 - 0.18, sub, ha='center', va='top',
                    fontsize=8, color='#777', style='italic', zorder=4)

    def ar(x1, y1, x2, y2, bi=False, clr='#999', lw=1.5):
        sty = '<|-|>' if bi else '-|>'
        a = FancyArrowPatch((x1, y1), (x2, y2),
                            arrowstyle=sty, mutation_scale=14,
                            color=clr, lw=lw, zorder=2)
        ax.add_patch(a)

    # Components
    bx(1.0, 2.5, 1.6, 0.8, M['dk'], 'Web Client', 'Browser / WS')
    bx(3.3, 2.5, 1.6, 0.8, M['red'], 'Nginx LB', 'hash(user_id)')
    bx(5.8, 4.0, 1.4, 0.6, M['dkred'], 'Gateway 1')
    bx(5.8, 2.5, 1.4, 0.6, M['dkred'], 'Gateway 2')
    bx(5.8, 1.0, 1.4, 0.6, M['dkred'], 'Gateway 3')
    bx(8.3, 2.5, 1.7, 0.8, M['dk'], 'Msg Service', 'Stateless / gRPC')
    bx(10.3, 3.3, 1.3, 0.6, M['red'], 'Redis', 'PTS (INCR)')
    bx(10.3, 1.7, 1.3, 0.6, M['blk'], 'Postgres', 'Messages')

    # Main flow arrows
    ar(1.8, 2.5, 2.5, 2.5)               # Client → Nginx
    ar(4.1, 2.75, 5.1, 3.85)             # Nginx → GW1
    ar(4.1, 2.5,  5.1, 2.5)              # Nginx → GW2
    ar(4.1, 2.25, 5.1, 1.15)             # Nginx → GW3
    ar(6.5, 3.7,  7.45, 2.8, clr='#bbb', lw=1.0)  # GW1 → Msg
    ar(6.5, 2.5,  7.45, 2.5)             # GW2 → Msg
    ar(6.5, 1.3,  7.45, 2.2, clr='#bbb', lw=1.0)  # GW3 → Msg
    ar(9.15, 2.75, 9.65, 3.15)           # Msg → Redis
    ar(9.15, 2.25, 9.65, 1.85)           # Msg → Postgres

    # gRPC mesh (vertical bidirectional)
    ar(6.65, 3.65, 6.65, 2.85, bi=True, clr=M['red'], lw=2.0)
    ar(6.65, 2.15, 6.65, 1.35, bi=True, clr=M['red'], lw=2.0)

    # Labels
    ax.text(2.15, 2.72, 'WS', fontsize=8, color='#888', ha='center')
    ax.text(7.0, 2.72, 'gRPC', fontsize=8, color='#888', ha='center')
    ax.text(7.15, 3.15, 'gRPC\nmesh', fontsize=8.5, color=M['red'],
            ha='left', va='center', fontweight='bold')
    ax.text(5.8, 4.55, 'Gateway Layer (Stateful)', fontsize=9.5,
            color=M['dkred'], ha='center', fontweight='bold')

    return tobuf(fig)


# ═══════════════════════════════════════════════════════════
#  CHART: EXPERIMENT 1 — Gateway Scaling
# ═══════════════════════════════════════════════════════════
def make_exp1():
    fig, ax = plt.subplots(figsize=(5.5, 4.2))
    fig.patch.set_facecolor('white')

    cfgs = ['1 Gateway\n(2000 conns)', '3 Gateways\n(~667 each)', '5 Gateways\n(~400 each)']
    mem = [55, 17, 10]
    colors = [M['red'], M['dkred'], M['dk']]

    bars = ax.bar(cfgs, mem, color=colors, width=0.55,
                  edgecolor='white', linewidth=2)

    for bar, v in zip(bars, mem):
        ax.text(bar.get_x() + bar.get_width()/2, v + 1.5,
                f'{v} MB', ha='center', va='bottom',
                fontweight='bold', fontsize=15, color='#333')

    ax.set_ylabel('Avg Heap per Gateway (MB)', fontsize=12, fontweight='bold')
    ax.set_title('Memory per Gateway at 2,000 Total Connections',
                 fontsize=13, fontweight='bold', pad=15, color='#222')
    ax.set_ylim(0, 68)
    ax.spines['bottom'].set_color('#ccc'); ax.spines['left'].set_color('#ccc')
    ax.yaxis.grid(True, alpha=0.25)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return tobuf(fig)


# ═══════════════════════════════════════════════════════════
#  CHART: EXPERIMENT 2 — Bottleneck
# ═══════════════════════════════════════════════════════════
def make_exp2():
    fig, ax = plt.subplots(figsize=(5.5, 3.6))
    fig.patch.set_facecolor('white')

    ops = ['Redis INCR\n(PTS Allocation)', 'Postgres INSERT\n(Message Store)']
    lat = [0.69, 3.10]
    colors = [M['red'], M['dk']]

    bars = ax.barh(ops, lat, color=colors, height=0.45,
                   edgecolor='white', linewidth=2)

    for bar, v in zip(bars, lat):
        ax.text(v + 0.1, bar.get_y() + bar.get_height()/2,
                f'{v:.2f} ms', va='center', fontweight='bold',
                fontsize=15, color='#333')

    ax.text(2.0, -0.55, '4.5\u00d7 difference', fontsize=13,
            fontweight='bold', color=M['red'], ha='center',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='#FFF0F0',
                      edgecolor=M['red'], lw=1.5))

    ax.set_xlabel('Average Latency (ms)', fontsize=12, fontweight='bold')
    ax.set_title('Write Path Latency Comparison',
                 fontsize=13, fontweight='bold', pad=15, color='#222')
    ax.set_xlim(0, 4.5)
    ax.spines['bottom'].set_color('#ccc'); ax.spines['left'].set_color('#ccc')
    ax.xaxis.grid(True, alpha=0.25)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return tobuf(fig)


# ═══════════════════════════════════════════════════════════
#  CHART: EXPERIMENT 3 — Failover
# ═══════════════════════════════════════════════════════════
def make_exp3():
    fig, ax = plt.subplots(figsize=(7, 4.2))
    fig.patch.set_facecolor('white')

    t   = [0,  3, 10, 20, 28, 30,  31, 33, 40, 50, 60, 80, 100, 130, 155]
    gw1 = [0, 22, 22, 22, 22, 22,   0,  0,  0,  0,  0,  0,   0,   0,   0]
    gw2 = [16,16, 16, 16, 16, 16,  16, 34, 34, 34, 18, 18,  18,  18,  18]
    gw3 = [12,12, 12, 12, 12, 12,  12, 26, 26, 26, 14, 14,  14,  14,  14]

    ax.plot(t, gw1, 'o-', color=M['red'], lw=2.5, ms=4, label='Gateway 1', zorder=3)
    ax.plot(t, gw2, 's-', color=M['dk'],  lw=2,   ms=4, label='Gateway 2', zorder=3)
    ax.plot(t, gw3, '^-', color=M['md'],  lw=2,   ms=4, label='Gateway 3', zorder=3)

    ax.axvline(x=31, color=M['red'], lw=2, ls='--', alpha=0.6, zorder=1)
    ax.axvspan(31, 33, alpha=0.08, color=M['red'], zorder=0)

    ax.annotate('GW1 killed', xy=(31, 1), xytext=(45, 5),
                fontsize=11, fontweight='bold', color=M['red'],
                arrowprops=dict(arrowstyle='->', color=M['red'], lw=1.5))
    ax.annotate('Connections\nredistributed', xy=(33, 34), xytext=(55, 37),
                fontsize=9.5, color=M['dk'],
                arrowprops=dict(arrowstyle='->', color='#999', lw=1.2))

    ax.set_xlabel('Time (seconds)', fontsize=12, fontweight='bold')
    ax.set_ylabel('Active Connections', fontsize=12, fontweight='bold')
    ax.set_title('Gateway Failover & Connection Redistribution',
                 fontsize=13, fontweight='bold', pad=15, color='#222')
    ax.legend(loc='upper right', fontsize=10, framealpha=0.95,
              edgecolor='#ddd')
    ax.set_ylim(-2, 42); ax.set_xlim(-5, 165)
    ax.spines['bottom'].set_color('#ccc'); ax.spines['left'].set_color('#ccc')
    ax.yaxis.grid(True, alpha=0.25)
    ax.set_axisbelow(True)
    fig.tight_layout()
    return tobuf(fig)


# ═══════════════════════════════════════════════════════════
#  BUILD SLIDES
# ═══════════════════════════════════════════════════════════

# ── 1. TITLE ──────────────────────────────────────────────
s = prs.slides.add_slide(BL); setbg(s, BLACK)
box(s, 0, 0, SW, 1.6, fill=RED)
box(s, 0, 1.6, SW, 0.15, fill=DKRED)

tx(s, 0, 2.3, SW, 1.0, 'TinyTelegram',
   sz=60, c=WHITE, b=True, al=PP_ALIGN.CENTER, fn='Segoe UI')
tx(s, 0, 3.3, SW, 0.5, 'A Distributed Instant Messaging System',
   sz=22, c=LTGRAY, al=PP_ALIGN.CENTER)

box(s, 4.5, 4.15, SW - 9, 0.015, fill=MDGRAY)

tx(s, 0, 4.5, SW, 0.4, 'CS6650  \u00b7  Building Scalable Distributed Systems',
   sz=16, c=MDGRAY, al=PP_ALIGN.CENTER)
tx(s, 0, 5.3, SW, 0.5, 'Fazheng Han   &   Emily Chen',
   sz=20, c=WHITE, b=True, al=PP_ALIGN.CENTER)
tx(s, 0, 5.95, SW, 0.4, 'Northeastern University  \u00b7  Spring 2026',
   sz=14, c=RED, al=PP_ALIGN.CENTER)

# ── 2. ARCHITECTURE ──────────────────────────────────────
s = prs.slides.add_slide(BL); setbg(s, WHITE)
hdr(s, 'System Architecture', 'Stateful/stateless split with Redis-based global ordering')
ftr(s, 2)

arch_buf = make_arch()
addimg(s, arch_buf, 0.3, 1.35, w=8.8)

concepts = [
    ('Consistent Hashing',  'Nginx hash(user_id) \u2192 O(1)\nrouting, minimal remapping'),
    ('PTS via Redis INCR',  'Atomic per-user sequence\nnumbers, no Raft needed'),
    ('gRPC Mesh Routing',   'Gateway-to-gateway delivery\nfor cross-node messages'),
    ('getDiff(pts) Sync',   'Incremental recovery on\nreconnect, zero message loss'),
]
cy = 1.55
for title, desc in concepts:
    box(s, 9.35, cy, 0.07, 0.85, fill=RED)
    tx(s, 9.6, cy, 3.3, 0.3, title, sz=13, b=True, c=DKGRAY)
    tx(s, 9.6, cy + 0.33, 3.3, 0.55, desc, sz=10, c=MDGRAY)
    cy += 1.25

# ── 3. EXPERIMENT 1 ──────────────────────────────────────
s = prs.slides.add_slide(BL); setbg(s, WHITE)
hdr(s, 'Experiment 1: Gateway Horizontal Scaling',
    '2,000 concurrent WebSocket connections across 1 / 3 / 5 gateway configurations')
ftr(s, 3)

addimg(s, make_exp1(), 0.4, 1.5, w=6.2)

box(s, 7.1, 1.5, 5.5, 5.3, fill=VLIGHT, rd=True)
ml(s, 7.4, 1.7, 5, 5, [
    ('Key Finding',                                 14, RED,    True,  PP_ALIGN.LEFT),
    ('',                                             6, WHITE,  False, PP_ALIGN.LEFT),
    ('Per-gateway memory scales linearly with',     13, DKGRAY, False, PP_ALIGN.LEFT),
    ('connections. Adding gateways distributes',    13, DKGRAY, False, PP_ALIGN.LEFT),
    ('load via consistent hashing.',                13, DKGRAY, False, PP_ALIGN.LEFT),
    ('',                                            10, WHITE,  False, PP_ALIGN.LEFT),
    ('Metrics @ 2,000 Total Connections',           13, RED,    True,  PP_ALIGN.LEFT),
    ('',                                             4, WHITE,  False, PP_ALIGN.LEFT),
    ('\u2022  1 GW:  55 MB/gw  (2000 conns each)', 12, DKGRAY, False, PP_ALIGN.LEFT),
    ('\u2022  3 GW:  17 MB/gw  (~667 conns each)', 12, DKGRAY, False, PP_ALIGN.LEFT),
    ('\u2022  5 GW:  10 MB/gw  (~400 conns each)', 12, DKGRAY, False, PP_ALIGN.LEFT),
    ('',                                            10, WHITE,  False, PP_ALIGN.LEFT),
    ('Total system memory ~50-55 MB regardless',    11, MDGRAY, False, PP_ALIGN.LEFT),
    ('of gateway count \u2014 load is distributed.',11, MDGRAY, False, PP_ALIGN.LEFT),
    ('',                                            14, WHITE,  False, PP_ALIGN.LEFT),
    ('Verdict:  \u2713 PASS',                       18, GREEN,  True,  PP_ALIGN.LEFT),
])

# ── 4. EXPERIMENT 2 ──────────────────────────────────────
s = prs.slides.add_slide(BL); setbg(s, WHITE)
hdr(s, 'Experiment 2: Write-Path Bottleneck Analysis',
    'Isolating Redis INCR vs Postgres INSERT latency in the persistence pipeline')
ftr(s, 4)

addimg(s, make_exp2(), 0.4, 1.7, w=6.2)

box(s, 7.1, 1.5, 5.5, 5.3, fill=VLIGHT, rd=True)
ml(s, 7.4, 1.7, 5, 5, [
    ('Key Finding',                                     14, RED,    True,  PP_ALIGN.LEFT),
    ('',                                                 6, WHITE,  False, PP_ALIGN.LEFT),
    ('Postgres INSERT is 4.5\u00d7 slower than',       13, DKGRAY, False, PP_ALIGN.LEFT),
    ('Redis INCR. Redis atomic operations are',        13, DKGRAY, False, PP_ALIGN.LEFT),
    ('NOT the bottleneck.',                             13, DKGRAY, True,  PP_ALIGN.LEFT),
    ('',                                                10, WHITE,  False, PP_ALIGN.LEFT),
    ('Latency Breakdown',                               13, RED,    True,  PP_ALIGN.LEFT),
    ('',                                                 4, WHITE,  False, PP_ALIGN.LEFT),
    ('\u2022  Redis INCR (PTS):      0.69 ms avg',     12, DKGRAY, False, PP_ALIGN.LEFT),
    ('\u2022  Postgres INSERT:       3.10 ms avg',     12, DKGRAY, False, PP_ALIGN.LEFT),
    ('',                                                10, WHITE,  False, PP_ALIGN.LEFT),
    ('Implication',                                     13, RED,    True,  PP_ALIGN.LEFT),
    ('',                                                 4, WHITE,  False, PP_ALIGN.LEFT),
    ('Redis INCR justified as PTS allocator \u2014',   11, MDGRAY, False, PP_ALIGN.LEFT),
    ('no need for heavier consensus (e.g. Raft).',     11, MDGRAY, False, PP_ALIGN.LEFT),
    ('To improve throughput: optimize PG writes',      11, MDGRAY, False, PP_ALIGN.LEFT),
    ('(batch inserts, connection pooling).',            11, MDGRAY, False, PP_ALIGN.LEFT),
    ('',                                                14, WHITE,  False, PP_ALIGN.LEFT),
    ('Verdict:  \u2713 PASS',                           18, GREEN,  True,  PP_ALIGN.LEFT),
])

# ── 5. EXPERIMENT 3 ──────────────────────────────────────
s = prs.slides.add_slide(BL); setbg(s, WHITE)
hdr(s, 'Experiment 3: Gateway Failover & Recovery',
    'Measuring reconnection behavior and message integrity when a gateway crashes')
ftr(s, 5)

addimg(s, make_exp3(), 0.2, 1.5, w=7.8)

box(s, 8.5, 1.5, 4.3, 5.3, fill=VLIGHT, rd=True)
ml(s, 8.75, 1.7, 3.8, 5, [
    ('Key Finding',                              14, RED,    True,  PP_ALIGN.LEFT),
    ('',                                          6, WHITE,  False, PP_ALIGN.LEFT),
    ('Zero message loss during',                 13, DKGRAY, False, PP_ALIGN.LEFT),
    ('failover. Connections auto-',              13, DKGRAY, False, PP_ALIGN.LEFT),
    ('redistribute to survivors.',               13, DKGRAY, False, PP_ALIGN.LEFT),
    ('',                                          8, WHITE,  False, PP_ALIGN.LEFT),
    ('Failover Sequence',                        12, RED,    True,  PP_ALIGN.LEFT),
    ('',                                          3, WHITE,  False, PP_ALIGN.LEFT),
    ('1. Baseline: 3 GWs, 50 conns',            10, DKGRAY, False, PP_ALIGN.LEFT),
    ('2. GW1 killed at t = 31s',                 10, DKGRAY, False, PP_ALIGN.LEFT),
    ('3. Clients reconnect to GW2/3',            10, DKGRAY, False, PP_ALIGN.LEFT),
    ('4. getDiff(pts) recovers msgs',            10, DKGRAY, False, PP_ALIGN.LEFT),
    ('',                                          8, WHITE,  False, PP_ALIGN.LEFT),
    ('Recovery Metrics',                          12, RED,    True,  PP_ALIGN.LEFT),
    ('',                                          3, WHITE,  False, PP_ALIGN.LEFT),
    ('\u2022  Drops: 22 (GW1 users)',             10, DKGRAY, False, PP_ALIGN.LEFT),
    ('\u2022  Recovery: < 3 seconds',             10, DKGRAY, False, PP_ALIGN.LEFT),
    ('\u2022  Message loss: 0',                   11, GREEN,  True,  PP_ALIGN.LEFT),
    ('',                                          12, WHITE,  False, PP_ALIGN.LEFT),
    ('Verdict:  \u2713 PASS',                     16, GREEN,  True,  PP_ALIGN.LEFT),
])

# ── 6. EXPERIMENT 4 ──────────────────────────────────────
s = prs.slides.add_slide(BL); setbg(s, WHITE)
hdr(s, 'Experiment 4: Consistency Validation',
    'Verifying PTS strict ordering under concurrent load with mid-test service crash')
ftr(s, 6)

# Big metric cards — row 1
cards1 = [
    ('218,059',  'Total Iterations',  DKGRAY),
    ('0',        'PTS Violations',    GREEN),
    ('1,211',    'Messages / sec',    DKGRAY),
]
for i, (num, label, c) in enumerate(cards1):
    cx = 1.2 + i * 4.1
    box(s, cx, 1.55, 3.2, 1.6, fill=VLIGHT, rd=True)
    tx(s, cx, 1.65, 3.2, 0.85, num,   sz=48, c=c,      b=True, al=PP_ALIGN.CENTER)
    tx(s, cx, 2.55, 3.2, 0.4,  label, sz=15, c=MDGRAY,          al=PP_ALIGN.CENTER)

# Separator
box(s, 1.2, 3.45, SW - 2.4, 0.012, fill=LTGRAY)

# Row 2 — test conditions
cards2 = [
    ('50 VUs',         'Concurrent Users'),
    ('2.5 min',        'Test Duration'),
    ('Service Crash',  'Mid-Test Restart'),
]
for i, (val, label) in enumerate(cards2):
    cx = 1.2 + i * 4.1
    box(s, cx, 3.7, 3.2, 1.15, fill=VLIGHT, rd=True)
    tx(s, cx, 3.78, 3.2, 0.6, val,   sz=26, c=RED,    b=True, al=PP_ALIGN.CENTER)
    tx(s, cx, 4.35, 3.2, 0.35, label, sz=13, c=MDGRAY,         al=PP_ALIGN.CENTER)

# Verdict banner
box(s, 2.2, 5.2, SW - 4.4, 1.2, fill=GREEN, rd=True)
tx(s, 2.2, 5.25, SW - 4.4, 0.65, 'VERDICT:  \u2713 PASS',
   sz=30, c=WHITE, b=True, al=PP_ALIGN.CENTER, anc=MSO_ANCHOR.MIDDLE)
tx(s, 2.2, 5.85, SW - 4.4, 0.45,
   'Redis INCR provides strict causal ordering without Raft consensus',
   sz=14, c=WHITE, al=PP_ALIGN.CENTER)

# ── 7. CLOUD DEPLOYMENT ──────────────────────────────────
s = prs.slides.add_slide(BL); setbg(s, WHITE)
hdr(s, 'Cloud Deployment: AWS Multi-AZ',
    'Production-grade infrastructure via AWS CDK  (Infrastructure as Code)')
ftr(s, 7)

# Left panel — AZ architecture
box(s, 0.8, 1.45, 5.6, 4.9, fill=VLIGHT, rd=True)

# AZ-a box
box(s, 1.1, 1.85, 2.5, 3.3, line=RED, rd=True)
tx(s, 1.1, 1.9, 2.5, 0.3, 'us-east-1a', sz=10, c=RED, b=True, al=PP_ALIGN.CENTER)
ml(s, 1.25, 2.25, 2.2, 2.8, [
    ('ECS Fargate',         11, DKGRAY, True,  PP_ALIGN.LEFT),
    ('  Gateway \u00d71',   10, MDGRAY, False, PP_ALIGN.LEFT),
    ('  MsgService \u00d71',10, MDGRAY, False, PP_ALIGN.LEFT),
    ('',                     6, WHITE,  False, PP_ALIGN.LEFT),
    ('Data (Primary)',       11, DKGRAY, True,  PP_ALIGN.LEFT),
    ('  RDS PostgreSQL',     10, MDGRAY, False, PP_ALIGN.LEFT),
    ('  ElastiCache Redis',  10, MDGRAY, False, PP_ALIGN.LEFT),
])

# AZ-b box
box(s, 3.8, 1.85, 2.5, 3.3, line=RED, rd=True)
tx(s, 3.8, 1.9, 2.5, 0.3, 'us-east-1b', sz=10, c=RED, b=True, al=PP_ALIGN.CENTER)
ml(s, 3.95, 2.25, 2.2, 2.8, [
    ('ECS Fargate',         11, DKGRAY, True,  PP_ALIGN.LEFT),
    ('  Gateway \u00d71',   10, MDGRAY, False, PP_ALIGN.LEFT),
    ('  MsgService \u00d71',10, MDGRAY, False, PP_ALIGN.LEFT),
    ('',                     6, WHITE,  False, PP_ALIGN.LEFT),
    ('Data (Standby)',       11, DKGRAY, True,  PP_ALIGN.LEFT),
    ('  RDS Read Replica',   10, MDGRAY, False, PP_ALIGN.LEFT),
    ('  ElastiCache Replica',10, MDGRAY, False, PP_ALIGN.LEFT),
])

# Edge row
tx(s, 0.8, 5.35, 5.6, 0.3, 'CloudFront + S3  (Static Web Client)',
   sz=11, c=RED, b=True, al=PP_ALIGN.CENTER)
tx(s, 0.8, 5.65, 5.6, 0.3, 'ALB  \u2192  ECS Service Discovery  \u2192  gRPC',
   sz=10, c=MDGRAY, al=PP_ALIGN.CENTER)
tx(s, 0.8, 5.95, 5.6, 0.3, 'AWS CDK (TypeScript) \u2014 all infra as code',
   sz=10, c=MDGRAY, al=PP_ALIGN.CENTER)

# Right panel — validation results
box(s, 6.8, 1.45, 5.7, 4.9, fill=VLIGHT, rd=True)
ml(s, 7.1, 1.6, 5.2, 4.7, [
    ('Cloud Validation Results',                   14, RED,    True,  PP_ALIGN.LEFT),
    ('',                                            8, WHITE,  False, PP_ALIGN.LEFT),
    ('Phase 1 \u2014 Data Plane',                  13, DKGRAY, True,  PP_ALIGN.LEFT),
    ('  \u2022  7,359 messages persisted (10 min)', 11, DKGRAY, False, PP_ALIGN.LEFT),
    ('  \u2022  0 duplicate PTS violations',        11, GREEN,  False, PP_ALIGN.LEFT),
    ('  \u2022  GW crash recovery: 3 seconds',      11, DKGRAY, False, PP_ALIGN.LEFT),
    ('',                                            8, WHITE,  False, PP_ALIGN.LEFT),
    ('Phase 2 \u2014 Compute (ECS Fargate)',       13, DKGRAY, True,  PP_ALIGN.LEFT),
    ('  \u2022  Multi-AZ deployment validated',     11, DKGRAY, False, PP_ALIGN.LEFT),
    ('  \u2022  Cross-AZ messaging confirmed',      11, DKGRAY, False, PP_ALIGN.LEFT),
    ('  \u2022  Service discovery operational',     11, DKGRAY, False, PP_ALIGN.LEFT),
    ('',                                            8, WHITE,  False, PP_ALIGN.LEFT),
    ('Phase 3 \u2014 Edge Layer',                  13, DKGRAY, True,  PP_ALIGN.LEFT),
    ('  \u2022  S3 + CloudFront static hosting',    11, DKGRAY, False, PP_ALIGN.LEFT),
    ('  \u2022  Live demo: 2 laptops, real-time',   11, DKGRAY, False, PP_ALIGN.LEFT),
    ('',                                            8, WHITE,  False, PP_ALIGN.LEFT),
    ('Technology Stack',                            13, RED,    True,  PP_ALIGN.LEFT),
    ('  AWS CDK (TS) \u00b7 ECS Fargate \u00b7 ALB',11,MDGRAY,False, PP_ALIGN.LEFT),
    ('  RDS PostgreSQL \u00b7 ElastiCache Redis',   11, MDGRAY, False, PP_ALIGN.LEFT),
    ('  S3 \u00b7 CloudFront \u00b7 Docker',        11, MDGRAY, False, PP_ALIGN.LEFT),
])

# ── 8. PHASE 1 — DATA PLANE VALIDATION ───────────────────
s = prs.slides.add_slide(BL); setbg(s, WHITE)
hdr(s, 'Phase 1: Cloud Data Plane Validation',
    'Local services \u2192 SSM tunnel \u2192 RDS + ElastiCache in us-east-1 Multi-AZ')
ftr(s, 8)

# Left — test (a): sustained load
box(s, 0.8, 1.45, 5.6, 2.6, fill=VLIGHT, rd=True)
tx(s, 1.05, 1.55, 5.1, 0.3, 'Test (a): 10-min Sustained Load', sz=14, c=RED, b=True)
ml(s, 1.05, 1.95, 5.1, 2.0, [
    ('5 VUs \u00d7 10 min via SSM bastion tunnel',    12, DKGRAY, False, PP_ALIGN.LEFT),
    ('',                                               4, WHITE,  False, PP_ALIGN.LEFT),
    ('\u2022  7,359 messages persisted',               12, DKGRAY, False, PP_ALIGN.LEFT),
    ('\u2022  0 duplicate receiver_pts',               12, GREEN,  True,  PP_ALIGN.LEFT),
    ('\u2022  0 duplicate sender_pts',                 12, GREEN,  True,  PP_ALIGN.LEFT),
    ('\u2022  ~12 msg/s (tunnel-bound, not in-VPC)',   11, MDGRAY, False, PP_ALIGN.LEFT),
    ('',                                               4, WHITE,  False, PP_ALIGN.LEFT),
    ('PTS gaps = CP behavior: writes rejected when',   10, MDGRAY, False, PP_ALIGN.LEFT),
    ('WAIT replica-ack timed out over tunnel RTT.',    10, MDGRAY, False, PP_ALIGN.LEFT),
])

# Left — test (b): kill test
box(s, 0.8, 4.3, 5.6, 2.35, fill=VLIGHT, rd=True)
tx(s, 1.05, 4.4, 5.1, 0.3, 'Test (b): Mid-Run Service Kill', sz=14, c=RED, b=True)
ml(s, 1.05, 4.78, 5.1, 1.8, [
    ('taskkill message-service at t = 60s, relaunch', 12, DKGRAY, False, PP_ALIGN.LEFT),
    ('',                                               4, WHITE,  False, PP_ALIGN.LEFT),
    ('\u2022  Downtime: 3 seconds (target \u2264 10s)',12, DKGRAY, False, PP_ALIGN.LEFT),
    ('\u2022  1,433 messages across 3 min',            12, DKGRAY, False, PP_ALIGN.LEFT),
    ('\u2022  0 PTS duplicates',                       12, GREEN,  True,  PP_ALIGN.LEFT),
    ('\u2022  5/5 VU sessions survived restart',       12, DKGRAY, False, PP_ALIGN.LEFT),
])

# Right — environment + timeline
box(s, 6.8, 1.45, 5.7, 3.2, fill=VLIGHT, rd=True)
tx(s, 7.05, 1.55, 5.2, 0.3, 'Test Environment', sz=14, c=RED, b=True)
env_items = [
    ('RDS PostgreSQL 16.3', 'db.m6g.large, Multi-AZ, 100 GB gp3'),
    ('ElastiCache Redis 7.1', 'cache.m6g.large, TLS + AUTH'),
    ('Connectivity', 'SSM Session Manager port-forward'),
    ('Services', 'Local Go binaries \u2192 cloud data'),
    ('Load Driver', 'k6 1.7.0, 5 VUs'),
]
ey = 1.95
for label, desc in env_items:
    tx(s, 7.05, ey, 2.3, 0.25, label, sz=10, c=DKGRAY, b=True)
    tx(s, 9.3,  ey, 3.0, 0.25, desc,  sz=10, c=MDGRAY)
    ey += 0.36

# Right bottom — verdict
box(s, 6.8, 4.9, 5.7, 1.75, fill=GREEN, rd=True)
ml(s, 7.05, 5.0, 5.2, 1.5, [
    ('\u2713  Phase 1 Exit Criteria: ALL PASSED',      16, WHITE, True,  PP_ALIGN.LEFT),
    ('',                                                6, WHITE, False, PP_ALIGN.LEFT),
    ('Data plane proves CP semantics hold on',         12, WHITE, False, PP_ALIGN.LEFT),
    ('real cloud infrastructure. PTS monotonicity',    12, WHITE, False, PP_ALIGN.LEFT),
    ('guaranteed: rejects writes it cannot replicate', 12, WHITE, False, PP_ALIGN.LEFT),
    ('rather than admitting out-of-order data.',       12, WHITE, False, PP_ALIGN.LEFT),
])

# ── 9. PHASE 2+3 — END-TO-END FARGATE DEMO ──────────────
s = prs.slides.add_slide(BL); setbg(s, WHITE)
hdr(s, 'Phase 2+3: End-to-End Fargate Demo',
    'CloudFront \u2192 ALB \u2192 ECS Fargate \u2192 RDS + ElastiCache  \u2014  two laptops, real-time messaging')
ftr(s, 9)

# Left — what we proved
box(s, 0.8, 1.45, 5.6, 2.4, fill=VLIGHT, rd=True)
tx(s, 1.05, 1.55, 5.1, 0.3, 'Cross-AZ Messaging Verified', sz=14, c=RED, b=True)
ml(s, 1.05, 1.95, 5.1, 1.8, [
    ('Two gateway tasks in two AZs (10.20.8.x,',      12, DKGRAY, False, PP_ALIGN.LEFT),
    ('10.20.7.x) \u2014 alice on GW-a, bob on GW-b',  12, DKGRAY, False, PP_ALIGN.LEFT),
    ('',                                                4, WHITE,  False, PP_ALIGN.LEFT),
    ('\u2022  alice \u2192 bob: gRPC peer call across AZs', 12, DKGRAY, False, PP_ALIGN.LEFT),
    ('\u2022  PTS monotonically incrementing',          12, GREEN,  True,  PP_ALIGN.LEFT),
    ('\u2022  "Sending\u2026" cleared on ack (both sides)',  12, DKGRAY, False, PP_ALIGN.LEFT),
])

# Left bottom — bugs found
box(s, 0.8, 4.1, 5.6, 2.55, fill=VLIGHT, rd=True)
tx(s, 1.05, 4.2, 5.1, 0.3, 'Bugs Found & Fixed in Production', sz=14, c=RED, b=True)
ml(s, 1.05, 4.58, 5.1, 2.0, [
    ('B1 \u2014 Gateway self-identity empty',          12, DKGRAY, True,  PP_ALIGN.LEFT),
    ('  ECS metadata not published to os.Environ;',    10, MDGRAY, False, PP_ALIGN.LEFT),
    ('  cross-GW delivery silently dropped.',           10, MDGRAY, False, PP_ALIGN.LEFT),
    ('',                                                3, WHITE,  False, PP_ALIGN.LEFT),
    ('B2 \u2014 Optimistic send bubble stuck',         12, DKGRAY, True,  PP_ALIGN.LEFT),
    ('  "Sending\u2026" never reconciled with ack;',   10, MDGRAY, False, PP_ALIGN.LEFT),
    ('  added FIFO pendingAcks queue.',                 10, MDGRAY, False, PP_ALIGN.LEFT),
    ('',                                                3, WHITE,  False, PP_ALIGN.LEFT),
    ('B3 \u2014 Missing logs:CreateLogGroup IAM',      12, DKGRAY, True,  PP_ALIGN.LEFT),
    ('  One-shot tasks failed; workaround: pre-create.',10, MDGRAY, False, PP_ALIGN.LEFT),
])

# Right — deployment topology
box(s, 6.8, 1.45, 5.7, 2.8, fill=VLIGHT, rd=True)
tx(s, 7.05, 1.55, 5.2, 0.3, 'Deployed Topology', sz=14, c=RED, b=True)
topo_items = [
    ('Web Client',      'S3 + CloudFront (public URL)'),
    ('Gateway',         '2 Fargate tasks, 2 AZs, behind ALB'),
    ('Message Service', '1 Fargate task, CloudMap DNS'),
    ('PostgreSQL',      'RDS db.m6g.large Multi-AZ'),
    ('Redis',           'ElastiCache cache.m6g.large Multi-AZ'),
    ('Stacks',          'VPC + Data + Compute + Edge (CDK)'),
]
ty = 1.95
for label, desc in topo_items:
    tx(s, 7.05, ty, 2.3, 0.25, label, sz=10, c=DKGRAY, b=True)
    tx(s, 9.3,  ty, 3.2, 0.25, desc,  sz=10, c=MDGRAY)
    ty += 0.36

# Right — live demo callout
box(s, 6.8, 4.5, 5.7, 2.15, fill=VLIGHT, rd=True)
tx(s, 7.05, 4.6, 5.2, 0.3, 'Live Demo Validation', sz=14, c=RED, b=True)
ml(s, 7.05, 5.0, 5.2, 1.5, [
    ('Two humans, two laptops, public URL:',           12, DKGRAY, False, PP_ALIGN.LEFT),
    ('d1ji1p758sdqkv.cloudfront.net',                  11, RED,    True,  PP_ALIGN.LEFT),
    ('',                                                4, WHITE,  False, PP_ALIGN.LEFT),
    ('\u2022  alice + bob exchanged messages',          11, DKGRAY, False, PP_ALIGN.LEFT),
    ('\u2022  Both sides rendered correctly',           11, DKGRAY, False, PP_ALIGN.LEFT),
    ('\u2022  PTS visible and monotonic',               11, GREEN,  True,  PP_ALIGN.LEFT),
    ('',                                                4, WHITE,  False, PP_ALIGN.LEFT),
    ('\u2713  Phase 2 + Phase 3 Exit: PASSED',         14, GREEN,  True,  PP_ALIGN.LEFT),
])

# ── 10. CONCLUSION ────────────────────────────────────────
s = prs.slides.add_slide(BL); setbg(s, BLACK)
box(s, 0, 0, SW, 0.95, fill=RED)
tx(s, 0, 0.05, SW, 0.85, 'Conclusion & Key Takeaways',
   sz=34, c=WHITE, b=True, al=PP_ALIGN.CENTER, anc=MSO_ANCHOR.MIDDLE)

results = [
    ('Exp 1', 'Horizontal Scaling',  '\u2713  Linear memory scaling validates consistent hashing'),
    ('Exp 2', 'Bottleneck Analysis', '\u2713  Redis INCR justified as PTS allocator (4.5\u00d7 faster)'),
    ('Exp 3', 'Failover Recovery',   '\u2713  Zero message loss with automatic reconnection'),
    ('Exp 4', 'Consistency',         '\u2713  218K iterations, zero PTS violations under failures'),
    ('Ph1',   'Data Plane',          '\u2713  7,359 msgs, 0 PTS duplicates, 3s crash recovery'),
    ('Ph2+3', 'Fargate Demo',        '\u2713  Cross-AZ messaging live on CloudFront + ECS'),
]

cy = 1.2
for tag, topic, finding in results:
    box(s, 1.0, cy, 1.5, 0.58, fill=DKRED, rd=True)
    tx(s, 1.0, cy, 1.5, 0.58, tag, sz=12, c=WHITE, b=True,
       al=PP_ALIGN.CENTER, anc=MSO_ANCHOR.MIDDLE)
    tx(s, 2.8, cy + 0.03, 2.5, 0.5, topic, sz=14, c=WHITE, b=True)
    tx(s, 5.6, cy + 0.03, 7, 0.5, finding, sz=13, c=LTGRAY)
    cy += 0.73

# DS concepts bar
box(s, 1.0, cy + 0.25, SW - 2, 0.012, fill=MDGRAY)
tx(s, 0, cy + 0.45, SW, 0.4,
   'Consistent Hashing  \u00b7  Redis INCR (PTS)  \u00b7  gRPC Mesh  '
   '\u00b7  getDiff Sync  \u00b7  CP Semantics',
   sz=13, c=RED, al=PP_ALIGN.CENTER)

# Thank you
tx(s, 0, SH - 1.8, SW, 0.7, 'Thank You!',
   sz=40, c=WHITE, b=True, al=PP_ALIGN.CENTER)
tx(s, 0, SH - 1.0, SW, 0.5, 'Questions?',
   sz=20, c=MDGRAY, al=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════
out_path = os.path.abspath(OUT)
prs.save(out_path)
print(f"Saved: {out_path}")
